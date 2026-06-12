"""
13_slides.py
------------
Expanded ~18-slide PowerPoint deck via python-pptx. Mohammad polishes in
PowerPoint after; the goal here is dense, defensible content rather than
visual polish.

Slide map:
   1  Title
   2  Research question & hypotheses
   3  Why this comparison matters (motivation)
   4  Study area
   5  Data: Sentinel-2 SR composites
   6  Land cover scheme (5 classes)
   7  Workflow overview (pipeline diagram-as-text)
   8  SLIC superpixel theory + chosen parameters
   9  SLIC segmentation result
  10  Per-segment features
  11  Training data + Random Forest configuration
  12  OBIA classified maps — 2016 & 2024
  13  Pixel-based pipeline + 3x3 majority filter
  14  Pixel-based classified maps — 2016 & 2024
  15  Side-by-side OBIA vs Pixel — 2024
  16  Validation design (stratified random, GE Pro)
  17  Accuracy assessment — headline table + confusion (2024)
  18  Per-class PA/UA — where each method fails
  19  Change detection — area chart + transitions
  20  Urban expansion 2016 -> 2024 map
  21  Discussion — why pixel beat OBIA at this scale
  22  Caveats & limitations
  23  Conclusions & future work
"""

from __future__ import annotations

from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

MAPS = Path(r"D:\Dehradun_OBIA\08_maps_final")
VAL = Path(r"D:\Dehradun_OBIA\06_validation")
CHG = Path(r"D:\Dehradun_OBIA\07_change_detection")
OBIA = Path(r"D:\Dehradun_OBIA\04_classification_OBIA")
PIX = Path(r"D:\Dehradun_OBIA\05_classification_pixel")
OUT = Path(r"D:\Dehradun_OBIA\09_report_slides")
OUT.mkdir(parents=True, exist_ok=True)

TITLE = "Object-Based vs Pixel-Based Classification for Dehradun Urban Expansion 2016-2024"
SUBTITLE = ("Mohammad Salman   |   M.Sc. Geo-informatics, IIRS Dehradun   |   "
            "Sentinel-2 SR + Python (rasterio, scikit-image, scikit-learn)")

NAVY = RGBColor(0x10, 0x2A, 0x43)
ACCENT = RGBColor(0xD7, 0x19, 0x1C)
MUTED = RGBColor(0x55, 0x55, 0x55)
GREEN = RGBColor(0x1A, 0x96, 0x50)


# -------------------------- low-level helpers --------------------------

def add_title_only(prs: Presentation, title: str) -> "Slide":
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
    slide.shapes.title.text = title
    p = slide.shapes.title.text_frame.paragraphs[0]
    p.font.color.rgb = NAVY
    p.font.size = Pt(26)
    p.font.bold = True
    return slide


def add_textbox(slide, left, top, width, height, lines, font_size=12,
                color=None, bold_first=False, line_spacing=1.05) -> None:
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.line_spacing = line_spacing
        if color:
            p.font.color.rgb = color
        if bold_first and i == 0:
            p.font.bold = True


def add_image(slide, path: Path, left, top, width=None, height=None) -> None:
    if not path.exists():
        add_textbox(slide, left, top, Inches(6), Inches(0.5),
                    [f"[missing: {path.name}]"], font_size=10, color=ACCENT)
        return
    if width and height:
        slide.shapes.add_picture(str(path), left, top, width=width, height=height)
    elif width:
        slide.shapes.add_picture(str(path), left, top, width=width)
    else:
        slide.shapes.add_picture(str(path), left, top)


def add_table(slide, left, top, width, height, data, header_bg=NAVY,
              header_fg=RGBColor(0xFF, 0xFF, 0xFF), font_size=10) -> None:
    """data: list of rows, first row is header."""
    rows = len(data)
    cols = len(data[0])
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    for j in range(cols):
        cell = table.cell(0, j)
        cell.text = str(data[0][j])
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(font_size)
            p.font.color.rgb = header_fg
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_bg
    for i in range(1, rows):
        for j in range(cols):
            cell = table.cell(i, j)
            cell.text = str(data[i][j])
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_size)


def add_footer(slide, page_n: int, total: int) -> None:
    add_textbox(slide, Inches(11.6), Inches(7.15), Inches(1.6), Inches(0.3),
                [f"{page_n} / {total}"], font_size=9, color=MUTED)


# -------------------------- build the deck ----------------------------

def build() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    total = 23  # final slide count

    # accuracy
    acc = pd.read_csv(VAL / "accuracy_summary.csv") if (VAL / "accuracy_summary.csv").exists() else None

    # --- 1 Title ---
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = TITLE
    slide.placeholders[1].text = SUBTITLE
    for p in slide.shapes.title.text_frame.paragraphs:
        p.font.size = Pt(30); p.font.color.rgb = NAVY; p.font.bold = True
    for p in slide.placeholders[1].text_frame.paragraphs:
        p.font.size = Pt(14); p.font.color.rgb = MUTED
    add_textbox(slide, Inches(0.6), Inches(6.6), Inches(12), Inches(0.5),
                ["Mini-project, IIRS Dehradun  |  June 2026"],
                font_size=12, color=MUTED)

    # --- 2 Research question & hypotheses ---
    slide = add_title_only(prs, "Research Question & Hypotheses")
    add_textbox(slide, Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.6),
                ["Research question:"], font_size=15, color=NAVY, bold_first=True)
    add_textbox(slide, Inches(0.7), Inches(1.85), Inches(12.0), Inches(1.4), [
        "Does Object-Based Image Analysis (OBIA) of Sentinel-2 produce a more accurate land-cover",
        "map of Dehradun than per-pixel classification — and how much new built-up area appeared",
        "between Nov 2016 (winter S2 composite) and Nov 2024?",
    ], font_size=13)
    add_textbox(slide, Inches(0.5), Inches(3.3), Inches(12.3), Inches(0.5),
                ["Hypotheses:"], font_size=15, color=NAVY, bold_first=True)
    add_textbox(slide, Inches(0.7), Inches(3.85), Inches(12.0), Inches(3.0), [
        "H1  OBIA will reduce salt-and-pepper noise compared to per-pixel Random Forest.",
        "H2  OBIA accuracy >= pixel-based accuracy because object-level features average",
        "      out within-segment noise and provide more stable training signatures.",
        "H3  Both methods will agree that Dehradun gained built-up area 2016 -> 2024,",
        "      primarily at the cost of cropland and dense vegetation.",
        "",
        "Spoiler (slide 17): H1 holds; H2 does NOT hold at the SLIC scale we chose; H3 holds.",
    ], font_size=13)
    add_footer(slide, 2, total)

    # --- 3 Why this comparison matters ---
    slide = add_title_only(prs, "Why Compare OBIA vs Pixel-Based?")
    add_textbox(slide, Inches(0.5), Inches(1.3), Inches(6.0), Inches(5.5), [
        "Pixel-based classification:",
        "  + Long-established (MLC, RF, SVM since 1990s).",
        "  + Captures fine spectral detail at sensor resolution.",
        "  - Salt-and-pepper noise in heterogeneous urban scenes.",
        "  - Ignores spatial context entirely.",
        "",
        "Object-Based Image Analysis (OBIA):",
        "  + Groups neighbouring spectrally-similar pixels into",
        "    'image objects' (segments) before classification.",
        "  + Each segment carries spatial + spectral + (optionally)",
        "    textural / shape descriptors.",
        "  + Visually smoother, polygon-ready output.",
        "  - Sensitive to segmentation parameters (scale, compactness).",
        "  - Mixed segments inherit a single class label.",
    ], font_size=12)
    add_textbox(slide, Inches(6.8), Inches(1.3), Inches(6.0), Inches(5.5), [
        "Open question for medium-resolution (10 m) data:",
        "  At 10 m Sentinel-2, individual urban objects (rooftops,",
        "  road segments, small fields) are 1-3 px wide. SLIC",
        "  segments of typical size (1000+ px) inevitably mix these.",
        "",
        "  Does the noise-reduction benefit of OBIA outweigh the",
        "  loss of within-segment spectral diversity?",
        "",
        "Existing literature is split:",
        "  - Blaschke (2010), Hossain & Chen (2019) — OBIA wins",
        "    on VHR imagery (< 2 m).",
        "  - Whyte et al. (2018), Ma et al. (2017) — comparable or",
        "    worse for medium-resolution.",
        "",
        "This study contributes one more data point in the second",
        "regime, applied to a rapidly urbanizing Himalayan foothill",
        "city.",
    ], font_size=12)
    add_footer(slide, 3, total)

    # --- 4 Study area ---
    slide = add_title_only(prs, "Study Area — Dehradun Valley, Uttarakhand")
    add_textbox(slide, Inches(0.5), Inches(1.3), Inches(5.0), Inches(5.5), [
        "Dehradun district capital, India",
        "  • Population: ~700k city, ~1.7M district (2024 est.)",
        "  • Capital of Uttarakhand, fastest-growing tier-2 city",
        "    in the Himalayan belt.",
        "  • Caught between Mussoorie ridge (N) and Shivalik (S).",
        "",
        "AOI extent:",
        "  • Bounding box ~58.6 x 33.7 km  (entire raster)",
        "  • Valid (in-AOI) area ~996 km^2  (after clip mask)",
        "  • CRS: EPSG:32644 (UTM zone 44N)",
        "  • Pixel size: 10 m",
        "",
        "Drivers of change 2016 -> 2024:",
        "  • IT/education sector expansion (Selaqui, Sahastradhara)",
        "  • Smart-City projects, peri-urban ribbon development",
        "  • Conversion of basmati cropland on the valley floor",
    ], font_size=12)
    add_image(slide, MAPS / "obia_map_2024.png", Inches(5.7), Inches(1.3), width=Inches(7.4))
    add_textbox(slide, Inches(5.7), Inches(6.3), Inches(7.4), Inches(0.4),
                ["OBIA-classified 2024 map showing the AOI shape."],
                font_size=10, color=MUTED)
    add_footer(slide, 4, total)

    # --- 5 Data ---
    slide = add_title_only(prs, "Data — Sentinel-2 Surface Reflectance Composites")
    add_textbox(slide, Inches(0.5), Inches(1.3), Inches(6.5), Inches(5.7), [
        "Two cloud-free winter composites from Google Earth Engine:",
        "",
        "  Dehradun_2016.tif  (Nov 2016 - Feb 2017)",
        "  Dehradun_2024.tif  (Nov 2024 - Feb 2025)",
        "",
        "Specifications (both scenes):",
        "  • 12 bands per scene, Float32, 10 m resolution",
        "  • CRS EPSG:32644, identical transform/shape — co-registered",
        "  • Raster size 3371 x 5864 px (~20M pixels, ~10M valid)",
        "  • Already clipped to Dehradun AOI in GEE",
        "",
        "Band order:",
        "  B2  B3  B4  B8  B5  B6  B7  B8A  B11  B12  NDVI  NDBI",
        "  Blue Green Red NIR Red-edge x3 (10m+resamp)  SWIR x2  indices",
        "",
        "Why winter window:",
        "  • Cloud-free in Indian Himalayan foothills.",
        "  • Dry season — clear distinction between irrigated cropland,",
        "    fallow bare land, and natural vegetation.",
    ], font_size=12)
    add_textbox(slide, Inches(7.2), Inches(1.3), Inches(6.0), Inches(5.7), [
        "Pre-processing audit (script 02):",
        "",
        "  Per-band statistics inside AOI revealed an unexpected issue:",
        "  B5, B6, B7, B8A, B11, B12 carry effectively NO signal",
        "  inside AOI — e.g. B11 std = ~22 DN on a ~1820 mean,",
        "  range 1744-1865 (~6% relative).",
        "",
        "  Likely cause: a band-reduction step in the GEE composite",
        "  export collapsed these bands. Validated by checking that",
        "  B2/B3/B4/B8/NDVI/NDBI all retain proper std (200-500 DN,",
        "  NDVI std ~0.22 — physically plausible).",
        "",
        "Decision:",
        "  Drop B5, B6, B7, B8A, B11, B12 from feature stack.",
        "  Final feature set (6): B2, B3, B4, B8, NDVI, NDBI.",
        "",
        "  Documented in 02_preprocessing/inside_aoi_stats.txt",
        "  for the Q&A reviewer.",
    ], font_size=12)
    add_footer(slide, 5, total)

    # --- 6 Class scheme ---
    slide = add_title_only(prs, "Land-Cover Class Scheme (5 classes)")
    add_table(
        slide, Inches(0.5), Inches(1.3), Inches(12.3), Inches(3.5),
        [
            ["ID", "Class", "Spectral signature (winter)", "Typical features in Dehradun"],
            ["1", "Built-up", "high B2/B3, NDBI > 0, NDVI < 0.3",
             "Old city, IT parks, ribbon settlements, roads"],
            ["2", "Dense Vegetation", "NDVI > 0.65, low SWIR",
             "Pine/oak forest on ridges, mature riparian belts"],
            ["3", "Mixed Veg / Cropland", "NDVI 0.30-0.55, mixed NDBI",
             "Basmati/wheat fields, orchards, sparse settlements"],
            ["4", "Bare / Open Land", "NDVI < 0.20, mid B8, mid NDBI",
             "Fallow fields, riverbeds, exposed soil, construction"],
            ["5", "Water", "very low NIR (B8), NDVI < 0.0",
             "Asan, Tons, Song rivers; small reservoirs"],
        ],
        font_size=11,
    )
    add_textbox(slide, Inches(0.5), Inches(5.1), Inches(12.3), Inches(2.0), [
        "Design rationale:",
        "  • Coarse enough for 10 m Sentinel-2 to discriminate reliably with only spectral features.",
        "  • Fine enough to track urban expansion (class 1) vs cropland conversion (class 3 -> 1, 4 -> 1).",
        "  • IGBP-compatible at the L1 hierarchy.",
        "  • Class 5 (water) is small in area at 10 m; we keep it for honesty even though both classifiers",
        "    will struggle with narrow rivers.",
    ], font_size=12)
    add_footer(slide, 6, total)

    # --- 7 Workflow ---
    slide = add_title_only(prs, "Workflow Overview")
    add_textbox(slide, Inches(0.5), Inches(1.3), Inches(6.2), Inches(5.8), [
        "OBIA pipeline:",
        "  (1) SLIC superpixels        scripts/03",
        "        n_segments = 8000, compactness = 10",
        "        input: B2, B3, B4, B8 (p1-p99 normalised)",
        "  (2) Per-segment mean features                     scripts/04",
        "        6 features: B2, B3, B4, B8, NDVI, NDBI",
        "  (3) Stratified random training        scripts/05",
        "        8 candidates per class per date (40/date)",
        "        Hint = simple NDVI/NDBI threshold rules",
        "        Truth = visual interpretation in GE Pro",
        "  (4) Random Forest classifier  scripts/07",
        "        300 trees, balanced class weights",
        "  (5) Predict 8000 segments -> rasterize @ 10 m",
        "",
        "Both pipelines share the SAME training data, the SAME",
        "feature set, and the SAME RF hyperparameters. Only the",
        "spatial unit of classification differs.",
    ], font_size=12)
    add_textbox(slide, Inches(7.0), Inches(1.3), Inches(6.0), Inches(5.8), [
        "Pixel-based pipeline:",
        "  (1) Training pixels = 3x3 window around training-segment centroids",
        "        ~9 training px / segment -> 360 px / date",
        "  (2) Random Forest (same config as OBIA)                scripts/08",
        "  (3) Predict every valid pixel (9.96 M / date)",
        "  (4) 3x3 majority filter to suppress salt-and-pepper",
        "  (5) Export classified GeoTIFF",
        "",
        "Validation (both methods):                                scripts/09-10",
        "  • 75 stratified random points per date (15 / class)",
        "  • Stratified on OBIA classified raster",
        "  • Visual interpretation in Google Earth Pro using",
        "    historical imagery (Nov-Feb of each year)",
        "  • Same 150 reference points used to score BOTH methods",
        "",
        "Change detection:                                         scripts/11",
        "  Cross-tabulate per method -> 5x5 transition matrix,",
        "  per-class area (ha), urbanization binary map.",
    ], font_size=12)
    add_footer(slide, 7, total)

    # --- 8 SLIC theory ---
    slide = add_title_only(prs, "SLIC Superpixel Segmentation — Theory")
    add_textbox(slide, Inches(0.5), Inches(1.3), Inches(6.0), Inches(5.8), [
        "Simple Linear Iterative Clustering (Achanta et al. 2012)",
        "",
        "  k-means in joint (x, y, spectral) space, restricted to a",
        "  local 2S x 2S neighbourhood per centroid where",
        "       S = sqrt(N_pixels / n_segments).",
        "",
        "Distance metric:",
        "    D = sqrt( (d_spectral / m)^2 + (d_spatial / S)^2 )",
        "  where m = compactness parameter.",
        "",
        "  - Low m   -> segments hug spectral edges (irregular)",
        "  - High m  -> segments stay close to a regular grid (compact)",
        "",
        "Computational complexity:",
        "  O(N) for the algorithm (vs O(N log N) for watershed,",
        "  O(N^2) for mean-shift).",
        "  Native scikit-image implementation, multi-threaded.",
    ], font_size=12)
    add_textbox(slide, Inches(6.8), Inches(1.3), Inches(6.0), Inches(5.8), [
        "Our parameter choices:",
        "",
        "  n_segments  = 8000",
        "    => with ~10M valid pixels, target ~1250 px / segment",
        "       ~ 12.5 ha per segment on a 10 m grid.",
        "",
        "  compactness = 10",
        "    => moderate — allow spectral edges to dominate",
        "       but maintain roughly square segments.",
        "",
        "  channels    = B2, B3, B4, B8",
        "    => visible RGB + NIR. NDVI/NDBI were NOT used for",
        "       segmentation (added later as features) so segments",
        "       are not biased toward our class definitions.",
        "",
        "  normalisation: per-channel p1-p99 min-max scaling.",
        "    Required because SLIC weights spatial vs spectral",
        "    distance using compactness; un-normalised DN scales",
        "    would dominate spatial distance.",
        "",
        "  mask = AOI valid mask  (start_label = 1)",
        "    Prevents segments from bleeding into outside-AOI.",
    ], font_size=12)
    add_footer(slide, 8, total)

    # --- 9 SLIC result ---
    slide = add_title_only(prs, "SLIC Segmentation Result")
    add_textbox(slide, Inches(0.5), Inches(1.3), Inches(4.5), Inches(5.5), [
        "Result — both dates:",
        "  segments inside AOI = 8000",
        "  median segment      = 1241 px  (~12.4 ha)",
        "  min  segment        = 748 px",
        "  max  segment        = ~2008 px",
        "  total runtime       = ~48 s / date",
        "",
        "Co-registration check:",
        "  Identical raster transform on both dates means",
        "  segment polygons are spatially equivalent, but",
        "  cluster IDs are not (SLIC is run independently",
        "  per date). This is appropriate because the actual",
        "  spectral content of a given location changes",
        "  between 2016 and 2024 — segmenting jointly would",
        "  bias the comparison.",
        "",
        "Computer / runtime:",
        "  Python 3.13, scikit-image 0.x, single workstation.",
    ], font_size=12)
    add_image(slide, MAPS / "obia_map_2024.png", Inches(5.3), Inches(1.3), width=Inches(7.7))
    add_textbox(slide, Inches(5.3), Inches(6.5), Inches(7.7), Inches(0.4),
                ["Each colour patch corresponds to one segment's predicted class — "
                 "the underlying segment boundaries are inherited from SLIC."],
                font_size=10, color=MUTED)
    add_footer(slide, 9, total)

    # --- 10 Per-segment features ---
    slide = add_title_only(prs, "Per-Segment Feature Extraction")
    add_textbox(slide, Inches(0.5), Inches(1.3), Inches(6.0), Inches(5.8), [
        "For each of the 8000 segments per date:",
        "",
        "  feature vector x  in  R^6",
        "    = [ mean(B2), mean(B3), mean(B4),",
        "        mean(B8), mean(NDVI), mean(NDBI) ]",
        "",
        "  ancillary fields (not used in classification):",
        "    n_pixels, centroid_row, centroid_col,",
        "    centroid_x, centroid_y (UTM 44N)",
        "    mean(B11)  -- retained for audit only",
        "",
        "Why mean rather than median:",
        "  • Cheap (O(N) via scipy.ndimage.mean).",
        "  • Sufficient with compact 12-ha segments; median was",
        "    tested on a sample and gave correlation > 0.999 vs mean.",
        "",
        "Why no texture features (GLCM, Gabor):",
        "  • Not in the brief; would add another paper's worth of",
        "    parameter tuning.",
        "  • At 10 m / 12 ha-segments, texture is dominated by",
        "    field boundaries -- captured implicitly by NDVI variance",
        "    across nearby segments.",
    ], font_size=12)
    add_textbox(slide, Inches(6.7), Inches(1.3), Inches(6.0), Inches(5.8), [
        "Feature distribution (2024, all 8000 segments):",
        "",
        "  mean_B2     range 134 -- 1756   std 242",
        "  mean_B3     range 251 -- 2066   std 265",
        "  mean_B4     range 169 -- 2338   std 344",
        "  mean_B8     range 929 -- 3593   std 237",
        "  mean_NDVI   range 0.005 -- 0.867 std 0.20",
        "  mean_NDBI   range -0.298 -- 0.474 std 0.06",
        "",
        "Important: at median 12 ha / segment, even a pure water",
        "body (Asan Barrage ~ 4 km^2) covers only ~30 segments,",
        "while small ponds (< 1 ha) are absorbed into surrounding",
        "vegetation segments. This explains the under-representation",
        "of class 5 (water) in both the candidate pool and the",
        "final maps.",
        "",
        "Persisted to:",
        "  04_classification_OBIA/segment_features_{year}.parquet",
        "  (also CSV for quick eyeballing)",
    ], font_size=12)
    add_footer(slide, 10, total)

    # --- 11 Training & RF ---
    slide = add_title_only(prs, "Training Data & Random Forest Configuration")
    add_textbox(slide, Inches(0.5), Inches(1.3), Inches(6.0), Inches(5.8), [
        "Training-segment selection (per date):",
        "  Step 1 — heuristic stratification (no learning yet):",
        "     • Built-up   NDBI > 0 AND NDVI < 0.30",
        "     • Dense Veg  NDVI > 0.65",
        "     • Mixed/Crop NDVI in [0.30, 0.55]",
        "     • Bare       NDVI < 0.20 AND NDBI in [-0.10, 0.10]",
        "                  AND B8 > 1500",
        "     • Water      B8 < 1500 AND NDVI < 0.25",
        "                  (fallback: 200 lowest-NIR segments)",
        "  Step 2 — sample 8 segments per class -> 40 candidates",
        "  Step 3 — visual confirmation in Google Earth Pro,",
        "          override the hint where wrong.",
        "",
        "  Final training distribution:",
        "     2016 -> {1:13, 2:9, 3:7, 4:7, 5:4}",
        "     2024 -> {1:12, 2:8, 3:7, 4:8, 5:5}",
        "  Reproducibility: numpy seed = 42.",
    ], font_size=11)
    add_textbox(slide, Inches(6.7), Inches(1.3), Inches(6.0), Inches(5.8), [
        "Random Forest hyperparameters (identical for both methods):",
        "",
        "  n_estimators        = 300",
        "  max_depth           = None  (grow fully)",
        "  min_samples_leaf    = 2",
        "  class_weight        = 'balanced'   (vital for class 5)",
        "  random_state        = 42",
        "  oob_score           = True",
        "",
        "Pixel-based training set:",
        "  3x3 window pulled at each training-segment centroid,",
        "  using the SAME class label -> ~360 px / date.",
        "  This deliberately keeps the methods comparable: both",
        "  classifiers see the same locations.",
        "",
        "OOB scores (small-sample, indicative only):",
        "  OBIA  2016 OOB = 0.70   2024 OOB = 0.75",
        "  Pixel 2016 OOB = ~0.80  2024 OOB = 0.81",
        "",
        "Real accuracy figures come from validation (slide 17).",
    ], font_size=11)
    add_footer(slide, 11, total)

    # --- 12 OBIA maps ---
    slide = add_title_only(prs, "OBIA Classified Maps — 2016 & 2024")
    add_image(slide, MAPS / "obia_map_2016.png", Inches(0.3), Inches(1.3), width=Inches(6.4))
    add_image(slide, MAPS / "obia_map_2024.png", Inches(6.7), Inches(1.3), width=Inches(6.4))
    add_textbox(slide, Inches(0.3), Inches(6.4), Inches(12.7), Inches(1.0), [
        "Note the visible expansion of red (Built-up) along the valley floor between Mussoorie (north) and",
        "Doiwala (south-east). Yellow (Bare/Open) patches in 2016 frequently transition to red in 2024 — the",
        "characteristic 'bare-land -> construction -> built-up' sequence on the urban fringe.",
    ], font_size=11, color=MUTED)
    add_footer(slide, 12, total)

    # --- 13 Pixel pipeline detail ---
    slide = add_title_only(prs, "Pixel-Based Pipeline & Post-Filtering")
    add_textbox(slide, Inches(0.5), Inches(1.3), Inches(6.0), Inches(5.8), [
        "Per-pixel feature vector  (R^6):",
        "  [ B2, B3, B4, B8, NDVI, NDBI ]  at every 10 m pixel.",
        "",
        "Prediction:",
        "  9 961 760 valid pixels per date scored in",
        "  ~1M-pixel chunks to bound RAM.",
        "  scikit-learn RF, n_jobs = -1 (all cores).",
        "",
        "Post-filtering — 3x3 majority filter:",
        "  • Standard salt-and-pepper suppression.",
        "  • For each pixel, replace with the modal class among",
        "    its 3x3 neighbourhood (zero/no-data ignored).",
        "  • Edge pixels keep their original class if all neighbours",
        "    are no-data.",
        "  • Implementation is fully vectorised over the 9-stack",
        "    of shifted neighbour arrays (script 08).",
        "",
        "Effect on area statistics:",
        "  Class areas changed by < 1% before/after filtering;",
        "  filter primarily removes isolated single pixels.",
    ], font_size=12)
    add_textbox(slide, Inches(6.7), Inches(1.3), Inches(6.0), Inches(5.8), [
        "Feature importances (Gini, RF, 2024):",
        "",
        "    NDVI    0.318",
        "    B2      0.185",
        "    B3      0.142",
        "    B4      0.128",
        "    B8      0.121",
        "    NDBI    0.105",
        "",
        "Same ranking as the OBIA RF — NDVI is the dominant",
        "discriminator. The leading role of B2 (blue) is notable:",
        "blue is sensitive to atmospheric haze and rooftop albedo,",
        "and the difference between concrete (high B2) and",
        "vegetation (low B2) is one of the strongest binary",
        "separators in the dataset.",
        "",
        "Note: NDBI ranks last despite being explicitly an urban",
        "index. Likely because at 10 m, NDBI is computed from",
        "the same NIR + (broken) SWIR information, with SWIR",
        "degraded; the index inherits NDVI's signal but with less",
        "discriminatory power.",
    ], font_size=12)
    add_footer(slide, 13, total)

    # --- 14 Pixel maps ---
    slide = add_title_only(prs, "Pixel-Based Classified Maps — 2016 & 2024")
    add_image(slide, MAPS / "pixel_map_2016.png", Inches(0.3), Inches(1.3), width=Inches(6.4))
    add_image(slide, MAPS / "pixel_map_2024.png", Inches(6.7), Inches(1.3), width=Inches(6.4))
    add_textbox(slide, Inches(0.3), Inches(6.4), Inches(12.7), Inches(1.0), [
        "After 3x3 majority filtering. Pixel-based maps preserve the fine 10 m structure of roads and",
        "individual built-up patches that SLIC segments dissolve into the surrounding land cover, particularly",
        "along the urban fringes south and west of the old city.",
    ], font_size=11, color=MUTED)
    add_footer(slide, 14, total)

    # --- 15 Side-by-side ---
    slide = add_title_only(prs, "OBIA vs Pixel-Based — 2024 Side-by-Side")
    add_image(slide, MAPS / "obia_vs_pixel_2024.png", Inches(0.2), Inches(1.3), width=Inches(12.9))
    add_textbox(slide, Inches(0.3), Inches(6.4), Inches(12.7), Inches(1.0), [
        "Pixel map (right) shows finer-grained classes; OBIA (left) shows polygon-like patches. Notice that",
        "OBIA covers the central built-up cluster as one large red patch while pixel preserves internal",
        "vegetation pockets (parks, riverine corridors).",
    ], font_size=11, color=MUTED)
    add_footer(slide, 15, total)

    # --- 16 Validation design ---
    slide = add_title_only(prs, "Validation Design")
    add_textbox(slide, Inches(0.5), Inches(1.3), Inches(6.0), Inches(5.8), [
        "Sample design — stratified random:",
        "  • Strata = OBIA classified classes (same strata used to",
        "    score BOTH methods so they are comparable).",
        "  • 15 points per stratum per date -> 75 per date -> 150 total.",
        "  • RNG seed = 2026 -> reproducible.",
        "  • Points exported as KML with class_obia and class_pixel",
        "    pre-attached, plus a UTM .gpkg for analysis.",
        "",
        "Reference labelling:",
        "  • Visual interpretation in Google Earth Pro using",
        "    historical imagery dated as close as possible to the",
        "    composite window (Nov-Feb of each year).",
        "  • Truth column class_true filled manually in CSV; the",
        "    KML pin name showed class_obia AND class_pixel so",
        "    bias from one method could be cross-checked.",
        "",
        "Why this design is defensible:",
        "  • Random within strata -> unbiased estimate per class.",
        "  • Stratification on OBIA ensures all classes get",
        "    representation (without stratification class 5 would",
        "    receive < 1 point in a uniform 75-point sample).",
    ], font_size=12)
    add_textbox(slide, Inches(6.7), Inches(1.3), Inches(6.0), Inches(5.8), [
        "Reference label distribution (true class):",
        "",
        "  2016:  18 BU  16 DV  23 MV   8 BL  10 W",
        "  2024:  19 BU  16 DV  22 MV  10 BL   8 W",
        "",
        "The true class differs from the OBIA stratum class for",
        "many points — that is exactly the signal the accuracy",
        "assessment is measuring.",
        "",
        "Limitations of the validation design:",
        "  • 75 points / date is small (per-class CI ~ +/-12 pp).",
        "  • Interpreter (Mohammad) is one person -- no",
        "    inter-rater reliability check.",
        "  • Historical GE imagery is not always exactly the same",
        "    season as the S2 composite; some judgement calls",
        "    on cropland phenology stage.",
        "",
        "  These limitations apply equally to both methods, so the",
        "  comparison itself is fair.",
    ], font_size=12)
    add_footer(slide, 16, total)

    # --- 17 Accuracy ---
    slide = add_title_only(prs, "Accuracy Assessment — Headline Results")
    rows = [["Method", "Year", "OA (%)", "Kappa", "n_val"]]
    if acc is not None:
        for _, r in acc.iterrows():
            rows.append([str(r["method"]), str(r["year"]),
                         f"{100*r['OA']:.1f}", f"{r['Kappa']:.2f}", str(int(r["n_val"]))])
    add_table(slide, Inches(0.5), Inches(1.3), Inches(6.0), Inches(2.5), rows, font_size=14)
    add_textbox(slide, Inches(0.5), Inches(4.1), Inches(6.0), Inches(3.0), [
        "Headline finding:",
        "  Pixel-based RF outperformed OBIA by 9-13 percentage",
        "  points OA and 0.10-0.15 in Kappa, on BOTH dates.",
        "",
        "H2 (OBIA accuracy >= pixel) is REJECTED at this segment",
        "scale.",
        "",
        "Both methods scored ~78-79% / kappa 0.72 for pixel and",
        "~65-69% / 0.57-0.62 for OBIA. The pixel-based method",
        "is consistent across dates; OBIA improved slightly from",
        "2016 to 2024 (richer urban texture to learn from).",
    ], font_size=12, bold_first=True)
    add_image(slide, VAL / "confusion_obia_2024.png", Inches(6.8), Inches(1.3), width=Inches(3.1))
    add_image(slide, VAL / "confusion_pixel_2024.png", Inches(10.0), Inches(1.3), width=Inches(3.1))
    add_image(slide, VAL / "confusion_obia_2016.png", Inches(6.8), Inches(4.5), width=Inches(3.1))
    add_image(slide, VAL / "confusion_pixel_2016.png", Inches(10.0), Inches(4.5), width=Inches(3.1))
    add_footer(slide, 17, total)

    # --- 18 Per-class --- (read accuracy_summary)
    slide = add_title_only(prs, "Per-Class Performance — Where Each Method Fails")
    if acc is not None:
        # build a per-class compact table
        # columns: class | OBIA PA / UA / F1 (2024) | Pixel PA / UA / F1 (2024)
        a24 = acc[acc["year"] == 2024]
        try:
            obia_row = a24[a24["method"] == "OBIA"].iloc[0]
            pix_row = a24[a24["method"] == "Pixel-based"].iloc[0]
            rows = [["Class", "OBIA  PA / UA / F1 (2024)", "Pixel  PA / UA / F1 (2024)"]]
            for c in ["Built-up", "Dense Veg", "Mixed Veg/Crop", "Bare/Open", "Water"]:
                rows.append([
                    c,
                    f"{obia_row[f'PA_{c}']*100:>5.1f}  /  {obia_row[f'UA_{c}']*100:>5.1f}  /  {obia_row[f'F1_{c}']:.2f}",
                    f"{pix_row[f'PA_{c}']*100:>5.1f}  /  {pix_row[f'UA_{c}']*100:>5.1f}  /  {pix_row[f'F1_{c}']:.2f}",
                ])
            add_table(slide, Inches(0.5), Inches(1.3), Inches(10.0), Inches(3.5), rows, font_size=11)
        except Exception:
            pass
    add_textbox(slide, Inches(0.5), Inches(5.0), Inches(12.3), Inches(2.5), [
        "Interpretation:",
        "  • Water — both methods perfect or near-perfect when they DO predict water; the bottleneck is",
        "    producer accuracy (omission). Narrow river pixels get absorbed into riparian vegetation.",
        "  • Mixed Veg/Crop — biggest omission class for OBIA. SLIC segments straddling field boundaries",
        "    get labeled by the dominant land cover, frequently 'Built-up' or 'Dense Veg', stripping cropland.",
        "  • Built-up — pixel-based catches small isolated rooftops that OBIA averages away into the segment",
        "    they sit in. This is the largest single accuracy gap between the methods.",
    ], font_size=11)
    add_footer(slide, 18, total)

    # --- 19 Change detection ---
    slide = add_title_only(prs, "Change Detection — Area & Class Transitions")
    add_image(slide, MAPS / "area_chart.png", Inches(0.3), Inches(1.3), width=Inches(13.0))
    add_textbox(slide, Inches(0.3), Inches(6.3), Inches(12.7), Inches(1.0), [
        "Annotated deltas in ha. Pixel-based (right) is the more conservative estimate of expansion;",
        "OBIA (left) over-estimates the BU class likely because some mixed cropland-settlement segments",
        "tip into the BU category. Both methods agree on direction for every class.",
    ], font_size=11, color=MUTED)
    add_footer(slide, 19, total)

    # --- 20 Urbanization map ---
    slide = add_title_only(prs, "Urban Expansion 2016 -> 2024")
    add_image(slide, MAPS / "change_map_pixel.png", Inches(0.3), Inches(1.3), width=Inches(8.5))
    add_textbox(slide, Inches(8.9), Inches(1.3), Inches(4.2), Inches(5.5), [
        "Pixel-based estimate:",
        "  Persistent built-up   ~ 7 750 ha",
        "  New built-up          ~ 8 020 ha",
        "  Net BU expansion      + 104%  (doubled)",
        "",
        "OBIA estimate (for comparison):",
        "  Persistent built-up   ~ 3 600 ha",
        "  New built-up          ~ 15 060 ha",
        "  Net BU expansion      + 384%",
        "",
        "Which to trust?",
        "  Pixel-based is the higher-accuracy classifier (slide 17),",
        "  so the pixel-based change number is the headline:",
        "",
        "         + ~80 km^2 of new built-up area in 8 years.",
        "",
        "  That is consistent with Dehradun's reported population",
        "  growth and IT-corridor expansion (Selaqui-Sahastradhara).",
        "  OBIA's higher number reflects its tendency to label mixed",
        "  fringe segments as Built-up.",
    ], font_size=12)
    add_footer(slide, 20, total)

    # --- 21 Discussion ---
    slide = add_title_only(prs, "Discussion — Why Pixel-Based Won at this Scale")
    add_textbox(slide, Inches(0.5), Inches(1.3), Inches(12.3), Inches(5.8), [
        "1.  Segment size dominates the result.",
        "    Our median segment is 12.4 ha. The characteristic length scale of land-cover change in",
        "    peri-urban Dehradun is ~ 0.1 -- 1 ha (single rooftops, narrow roads, small fields). When",
        "    a 12-ha SLIC segment straddles a boundary between built-up and cropland, the within-",
        "    segment average pulls toward whichever cover dominates, and the minority class is lost.",
        "",
        "2.  Pixel-RF is tolerant of small training samples.",
        "    The 9x training-pixel multiplier (3x3 window) gives the pixel RF 360 training samples vs",
        "    only 40 for OBIA. This effectively closes the bias gap from limited fieldwork.",
        "",
        "3.  Spectral confusion classes resolve differently.",
        "    Bare-soil construction sites have spectral signatures close to built-up rooftops. At the",
        "    pixel level the RF learns this nuance from 360 samples; at the segment level, with only 7-8",
        "    samples per class, the boundary is fuzzier.",
        "",
        "4.  3x3 majority filter is enough.",
        "    The principal critique of pixel-based classification -- salt-and-pepper -- is largely cured",
        "    by a single 3x3 mode filter, without sacrificing boundary fidelity.",
        "",
        "5.  This does NOT mean OBIA is bad. It means OBIA at this segment scale is not appropriate.",
        "    Two natural fixes: (a) shrink segments (n_segments = 30 000+), (b) multi-scale OBIA",
        "    (e.g. Berkeley's multiresolution segmentation at multiple scales, or hierarchical SLIC).",
    ], font_size=12)
    add_footer(slide, 21, total)

    # --- 22 Caveats ---
    slide = add_title_only(prs, "Caveats & Limitations")
    add_textbox(slide, Inches(0.5), Inches(1.3), Inches(12.3), Inches(5.8), [
        "Data quality:",
        "  • Red-edge (B5, B6, B7, B8A) and SWIR (B11, B12) bands inside the AOI were effectively",
        "    constant in both composites -- almost certainly a GEE export-side compositing artefact.",
        "    Final feature set restricted to B2, B3, B4, B8, NDVI, NDBI. Re-running with a cleaner",
        "    composite (raw L2A SR product) would test whether adding SWIR closes the OBIA accuracy",
        "    gap, since NDBI would gain real signal.",
        "",
        "Methodological:",
        "  • Only one segmentation scale tested. Sensitivity analysis (n_segments in {2000, 8000, 30000})",
        "    would reveal whether the OBIA accuracy gap is fundamental or just a scale choice.",
        "  • Training sample n = 40 per date for OBIA is small. RF generalisation may have saturated.",
        "  • No texture features used (GLCM, Gabor). Adding them might help OBIA more than pixel,",
        "    since OBIA can pool texture statistics over the whole segment.",
        "",
        "Validation:",
        "  • 75 reference points per date; per-class confidence intervals roughly +/- 12 pp at p=0.05.",
        "  • Single interpreter -- no inter-rater reliability quantified.",
        "  • Historical imagery in GE Pro is not always the same season as the S2 composite.",
        "",
        "Temporal:",
        "  • Two snapshots; no information about transition trajectories (e.g. when did each new",
        "    built-up pixel transition?). A dense time series would address this.",
    ], font_size=12)
    add_footer(slide, 22, total)

    # --- 23 Conclusions ---
    slide = add_title_only(prs, "Conclusions & Future Work")
    add_textbox(slide, Inches(0.5), Inches(1.3), Inches(12.3), Inches(5.8), [
        "Conclusions:",
        "  1.  For 10 m Sentinel-2 in a mixed peri-urban setting, a per-pixel Random Forest with a",
        "      simple 3x3 majority filter outperforms a single-scale SLIC-OBIA Random Forest using",
        "      identical training data, features and hyperparameters.  (OA 78.7% vs 65.3-69.3%).",
        "  2.  Both methods agree that Dehradun gained substantial built-up area 2016 -> 2024.",
        "      The defensible (pixel-based) figure is ~ 80 km^2 of NEW built-up area over 8 years,",
        "      mostly converted from mixed cropland and bare fallow on the valley floor.",
        "  3.  Dense vegetation lost ~ 130-160 km^2 over the same period -- the largest single",
        "      land-cover signal in the dataset.",
        "",
        "Future work:",
        "  • Multi-scale OBIA -- test n_segments in {2k, 8k, 30k, 80k} to find the scale where",
        "    OBIA matches pixel-based, or train an ensemble across scales.",
        "  • Re-export Sentinel-2 composite preserving red-edge and SWIR bands; re-classify with",
        "    full 12-feature stack.",
        "  • Add texture features (GLCM mean/variance) at segment level.",
        "  • Inter-rater reliability: independent reference labelling by 2-3 interpreters.",
        "  • Extend time series to annual scale 2016-2024 (9 dates) to capture transition dynamics.",
        "  • Compare against deep-learning baselines (U-Net, DeepLabv3+) at 10 m.",
        "",
        "All code, intermediate artefacts, and this deck are reproducible from",
        "  D:/Dehradun_OBIA/scripts/01_inspect_rasters.py ... 13_slides.py",
    ], font_size=11)
    add_footer(slide, 23, total)

    pptx_out = OUT / "Dehradun_OBIA_vs_Pixel.pptx"
    try:
        prs.save(pptx_out)
    except PermissionError:
        # PPTX is open in PowerPoint; save under a fallback name so we never
        # silently fail the run.
        pptx_out = OUT / "Dehradun_OBIA_vs_Pixel_v2.pptx"
        prs.save(pptx_out)
        print(f"NOTE: original was locked (open in PowerPoint?). Saved as {pptx_out.name}")
    print(f"Saved {pptx_out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
